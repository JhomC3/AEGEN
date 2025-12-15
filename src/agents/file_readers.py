# src/agents/file_readers.py
import csv
from collections.abc import Callable

import docx
import openpyxl
import pptx
import pypdf

READER_REGISTRY: dict[str, Callable[[str], str]] = {}


def register_reader(*extensions: str):
    """Decorador para registrar lectores por extensión."""
    def decorator(func: Callable[[str], str]) -> Callable[[str], str]:
        for ext in extensions:
            READER_REGISTRY[ext.lower()] = func
        return func
    return decorator


@register_reader(".txt", ".md")
def read_txt(file_path: str) -> str:
    """Lee archivos de texto plano."""
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        content = f.read()
    if not content.strip():
        raise ValueError("Text file is empty")
    return content


@register_reader(".pdf")
def read_pdf(file_path: str) -> str:
    """Extrae texto de PDF."""
    content = []
    try:
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    content.append(text)
    except Exception as e:
        raise ValueError(f"Error reading PDF: {str(e)}") from e

    if not content:
        raise ValueError("PDF appears to be empty")
    return "\n\n".join(content)


@register_reader(".docx")
def read_docx(file_path: str) -> str:
    """Extrae texto de DOCX."""
    try:
        document = docx.Document(file_path)
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        if not paragraphs:
            raise ValueError("DOCX is empty")
        return "\n".join(paragraphs)
    except Exception as e:
        raise ValueError(f"Error reading DOCX: {str(e)}") from e


@register_reader(".pptx")
def read_pptx(file_path: str) -> str:
    """Extrae texto de PPTX."""
    try:
        presentation = pptx.Presentation(file_path)
        content = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
        if not content:
            raise ValueError("PPTX is empty")
        return "\n".join(content)
    except Exception as e:
        raise ValueError(f"Error reading PPTX: {str(e)}") from e


@register_reader(".xlsx")
def read_xlsx(file_path: str) -> str:
    """Extrae texto de XLSX."""
    try:
        workbook = openpyxl.load_workbook(file_path, read_only=True)
        content = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                row_data = [str(cell.value) for cell in row if cell.value is not None]
                if row_data:
                    content.append(", ".join(row_data))
        if not content:
            raise ValueError("XLSX is empty")
        return "\n".join(content)
    except Exception as e:
        raise ValueError(f"Error reading XLSX: {str(e)}") from e


@register_reader(".csv")
def read_csv(file_path: str) -> str:
    """Extrae texto de CSV."""
    try:
        content = []
        with open(file_path, encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                if row:
                    content.append(", ".join(row))
        if not content:
            raise ValueError("CSV is empty")
        return "\n".join(content)
    except Exception as e:
        raise ValueError(f"Error reading CSV: {str(e)}") from e
