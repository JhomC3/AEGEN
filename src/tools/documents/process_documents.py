import asyncio
import csv
import logging
from collections.abc import Callable
from pathlib import Path

import docx
import openpyxl
import pptx
import pypdf
from core.schemas import DocumentContent, DocumentError

READER_REGISTRY: dict[str, Callable[[str], str]] = {}


def register_reader(*extensions: str):
    """
    Decorador para registrar una funcion como lectora de una o más extensiones.
    """

    def decorador(func: Callable[[str], str]) -> Callable[[str], str]:
        for ext in extensions:
            READER_REGISTRY[ext] = func
        return func

    return decorador


@register_reader(".txt", ".md")
def read_txt(file_path: str) -> str:
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        return f.read()


@register_reader(".pdf")
def read_pdf(file_path: str) -> str:
    """Extrae el texto de un archivo PDF."""
    content = []
    with open(file_path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            content.append(page.extract_text() or "")
    return "\n\n".join(content)


@register_reader(".docx")
def read_docx(file_path: str) -> str:
    """Extrae el texto de un archivo DOCX."""
    document = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in document.paragraphs])


@register_reader(".pptx")
def read_pptx(file_path: str) -> str:
    """Extrae el texto de un archivo PPTX."""
    presentation = pptx.Presentation(file_path)
    content = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
    return "\n".join(content)


@register_reader(".xlsx")
def read_xlsx(file_path: str) -> str:
    """Extrae el texto de un archivo XLSX."""
    workbook = openpyxl.load_workbook(file_path, read_only=True)
    content = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            row_content = [str(cell.value) for cell in row if cell.value is not None]
            if row_content:
                content.append(", ".join(row_content))
    return "\n".join(content)


@register_reader(".csv")
def read_csv(file_path: str) -> str:
    """Extrae el texto de un archivo CSV."""
    content = []
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            content.append(", ".join(row))
    return "\n".join(content)


class DocumentProcessor:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.readers = READER_REGISTRY
        self.logger.info(
            f"DocumentProcessor initialized with readers: {list(self.readers.keys())}"
        )

    async def process_file(
        self,
        file_path: str,
        file_name: str,
    ) -> DocumentContent | DocumentError:
        """Procesa un archivo de documento para extraer su contenido."""
        self.logger.info(f"Iniciando procesamiento del archivo: {file_name}")
        file_extension = Path(file_name).suffix.lower()
        reader_func = self.readers.get(file_extension)
        if not reader_func:
            msg = f"Formato de archivo no soportado: {file_extension}"
            self.logger.warning(msg)
            return DocumentError(
                file_name=file_name,
                message=msg,
                extension=file_extension,
            )
        try:
            self.logger.info(f"Usando lector para {file_extension}")
            content = await asyncio.to_thread(reader_func, file_path)
            self.logger.info(f"Exito al procesar el archivo: {file_name}")
            return DocumentContent(
                file_name=file_name,
                content=content,
                extension=file_extension,
            )
        except Exception as e:
            msg = f"Error al procesar el archivo: {file_name} - {e}"
            self.logger.error(msg, exc_info=True)
            return DocumentError(
                file_name=file_name,
                message=msg,
                extension=file_extension,
            )
