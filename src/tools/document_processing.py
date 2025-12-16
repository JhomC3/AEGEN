import asyncio
import csv
import logging
from collections.abc import Callable
from pathlib import Path

import docx
import openpyxl
import pptx
import pypdf
from langchain_core.tools import tool

logger = logging.getLogger(__name__)

# --- Registry for Document Readers ---
READER_REGISTRY: dict[str, Callable[[str], str]] = {}


def register_reader(*extensions: str):
    """
    Decorador para registrar una función como lectora de una o más extensiones.
    Esto permite añadir fácilmente soporte para nuevos tipos de archivo.
    """

    def decorator(func: Callable[[str], str]) -> Callable[[str], str]:
        for ext in extensions:
            READER_REGISTRY[ext.lower()] = func
        return func

    return decorator


# --- Reader Implementations ---


@register_reader(".txt", ".md")
def read_txt(file_path: str) -> str:
    """Lee archivos de texto plano."""
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        return f.read()


@register_reader(".pdf")
def read_pdf(file_path: str) -> str:
    """Extrae el texto de un archivo PDF."""
    content = []
    with open(file_path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            content.append(page.extract_text() or "")
    return "\n\n".join(content)


@register_reader(".docx")
def read_docx(file_path: str) -> str:
    """Extrae el texto de un archivo DOCX."""
    document = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in document.paragraphs])


@register_reader(".pptx")
def read_pptx(file_path: str) -> str:
    """Extrae el texto de un archivo PPTX."""
    presentation = pptx.Presentation(file_path)
    content = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
    return "\n".join(content)


@register_reader(".xlsx")
def read_xlsx(file_path: str) -> str:
    """Extrae el texto de un archivo XLSX."""
    workbook = openpyxl.load_workbook(file_path, read_only=True)
    content = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            row_content = [str(cell.value) for cell in row if cell.value is not None]
            if row_content:
                content.append(", ".join(row_content))
    return "\n".join(content)


@register_reader(".csv")
def read_csv(file_path: str) -> str:
    """Extrae el texto de un archivo CSV."""
    content = []
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            content.append(", ".join(row))
    return "\n".join(content)


# --- Main Processing Tool ---

@tool
async def process_document(file_path: str, file_name: str) -> dict:
    """
    Procesa un archivo de documento para extraer su contenido de forma asíncrona.
    Selecciona el lector adecuado según la extensión del archivo.
    """
    logger.info(f"Processing document: {file_name} at {file_path}")
    file_extension = Path(file_name).suffix.lower()
    reader_func = READER_REGISTRY.get(file_extension)

    if not reader_func:
        error_message = f"El procesamiento para archivos '{file_extension}' no está implementado."
        logger.warning(error_message)
        return {"error": error_message}

    try:
        logger.info(f"Using reader for '{file_extension}'")
        # Ejecuta la función de lectura (que es I/O-bound) en un hilo separado
        content = await asyncio.to_thread(reader_func, file_path)
        logger.info(f"Successfully processed document: {file_name}")
        return {"content": content}
    except Exception as e:
        error_message = f"Error processing document {file_path}: {e}"
        logger.error(error_message, exc_info=True)
        return {"error": str(e)}
